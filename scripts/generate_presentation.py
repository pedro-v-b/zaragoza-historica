"""Genera la presentacion del TFG Zaragoza Historica (10 diapositivas).

Requisitos cumplidos:
- 10 slides (~1 min cada una = 10 min total)
- Texto mínimo (maximo ~6 palabras clave por slide)
- Estructura QUÉ / CÓMO / POR QUÉ equilibrada
- Diseño limpio con paleta del proyecto (azul #5A8A9F + naranja #F59435)
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Paleta del proyecto
AZUL = RGBColor(0x5A, 0x8A, 0x9F)
AZUL_CLARO = RGBColor(0x7C, 0xA5, 0xB8)
NARANJA = RGBColor(0xF5, 0x94, 0x35)
FONDO = RGBColor(0xFF, 0xFF, 0xFF)
TEXTO = RGBColor(0x2C, 0x3E, 0x50)
TEXTO_MUTE = RGBColor(0x7F, 0x8C, 0x8D)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)

LOGO_PATH = Path(r"C:\Users\pvial\Desktop\TFG DAM v2\frontend\public\logo.png")
OUT_PATH = Path(r"C:\Users\pvial\Desktop\TFG DAM v2\Zaragoza_Historica_TFG_v2.pptx")

# Widescreen 16:9
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

BLANK = prs.slide_layouts[6]


def set_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text(slide, left, top, width, height, text, size, color=TEXTO, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return tb


def add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.6), color=NARANJA):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_circle_stat(slide, cx, cy, size, number, label):
    """Círculo naranja con número grande + etiqueta debajo."""
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - size / 2, cy - size / 2, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = NARANJA
    circle.line.fill.background()
    tf = circle.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = number
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.color.rgb = BLANCO
    r.font.name = "Calibri"
    # Etiqueta
    add_text(slide, cx - Inches(1.5), cy + size / 2 + Inches(0.1), Inches(3), Inches(0.4),
             label, 14, color=TEXTO, align=PP_ALIGN.CENTER)


def footer(slide, page_num, total=10):
    add_text(slide, Inches(0.5), Inches(7.0), Inches(6), Inches(0.35),
             "Zaragoza Histórica  |  TFG DAM", 10, color=TEXTO_MUTE)
    add_text(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.35),
             f"{page_num} / {total}", 10, color=TEXTO_MUTE, align=PP_ALIGN.RIGHT)


def screenshot_hint(slide, text, left=Inches(0.5), top=Inches(6.55), width=Inches(12.3)):
    """Añade una nota con la recomendación de captura (editable/borrable)."""
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = f"({text})"
    r.font.size = Pt(10)
    r.font.italic = True
    r.font.color.rgb = NARANJA
    r.font.name = "Calibri"


# ========== SLIDE 1: PORTADA ==========
s = prs.slides.add_slide(BLANK)
set_bg(s, FONDO)
# Banda lateral izquierda azul
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(4.5), SLIDE_H)
band.fill.solid()
band.fill.fore_color.rgb = AZUL
band.line.fill.background()
# Logo en la banda
if LOGO_PATH.exists():
    s.shapes.add_picture(str(LOGO_PATH), Inches(1.25), Inches(2.5), width=Inches(2), height=Inches(2))
# Título y subtítulo
add_text(s, Inches(5.2), Inches(2.3), Inches(7.5), Inches(1.2),
         "Zaragoza Histórica", 48, color=AZUL, bold=True)
add_text(s, Inches(5.2), Inches(3.3), Inches(7.5), Inches(0.6),
         "Fotografía histórica geolocalizada", 22, color=NARANJA)
# Separador
sep = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(4.1), Inches(1.5), Emu(15000))
sep.fill.solid()
sep.fill.fore_color.rgb = NARANJA
sep.line.fill.background()
add_text(s, Inches(5.2), Inches(4.4), Inches(7.5), Inches(0.5),
         "Trabajo Fin de Grado  ·  DAM", 16, color=TEXTO)
add_text(s, Inches(5.2), Inches(5.0), Inches(7.5), Inches(0.5),
         "Pedro", 14, color=TEXTO_MUTE)
add_text(s, Inches(5.2), Inches(5.4), Inches(7.5), Inches(0.5),
         "2026", 14, color=TEXTO_MUTE)
screenshot_hint(s, "Captura recomendada: vista general del mapa con marcadores y clusters visibles, como hero background")


# ========== SLIDE 2: EL QUÉ (Problema) ==========
s = prs.slides.add_slide(BLANK)
set_bg(s, FONDO)
add_accent_bar(s, Inches(0.5), Inches(0.6))
add_text(s, Inches(0.75), Inches(0.55), Inches(12), Inches(0.7),
         "El problema", 32, color=AZUL, bold=True)

# Tres puntos clave
points = [
    ("📚", "Archivos dispersos", "AMZ, Flickr, fondos digitales"),
    ("📍", "Sin contexto espacial", "¿Dónde fue tomada cada foto?"),
    ("🔍", "Difícil exploración", "No hay búsqueda visual unificada"),
]
for i, (icon, title, desc) in enumerate(points):
    x = Inches(0.8 + i * 4.2)
    y = Inches(2.5)
    # Caja
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(3))
    box.fill.solid()
    box.fill.fore_color.rgb = BLANCO
    box.line.color.rgb = AZUL_CLARO
    box.line.width = Pt(1)
    # Icono
    add_text(s, x, y + Inches(0.3), Inches(3.8), Inches(0.8), icon, 44, align=PP_ALIGN.CENTER)
    # Título
    add_text(s, x + Inches(0.2), y + Inches(1.4), Inches(3.4), Inches(0.5),
             title, 18, color=AZUL, bold=True, align=PP_ALIGN.CENTER)
    # Descripción
    add_text(s, x + Inches(0.2), y + Inches(2.0), Inches(3.4), Inches(0.6),
             desc, 13, color=TEXTO_MUTE, align=PP_ALIGN.CENTER)
screenshot_hint(s, "Captura recomendada: collage de 3-4 webs de archivos fotográficos dispersos (AMZ, Flickr, DARA) mostrando la fragmentación")
footer(s, 2)


# ========== SLIDE 3: LA SOLUCIÓN ==========
s = prs.slides.add_slide(BLANK)
set_bg(s, FONDO)
add_accent_bar(s, Inches(0.5), Inches(0.6))
add_text(s, Inches(0.75), Inches(0.55), Inches(12), Inches(0.7),
         "La solución", 32, color=AZUL, bold=True)
add_text(s, Inches(0.75), Inches(1.3), Inches(12), Inches(0.5),
         "Mapa interactivo web con exploración visual", 18, color=TEXTO_MUTE)

# Tres círculos de estadísticas
add_circle_stat(s, Inches(3.0), Inches(4.5), Inches(2.2), "5.214", "fotografías")
add_circle_stat(s, Inches(6.67), Inches(4.5), Inches(2.2), "35K", "edificios")
add_circle_stat(s, Inches(10.33), Inches(4.5), Inches(2.2), "11", "capas mapa")
screenshot_hint(s, "Captura recomendada: vista general del mapa con clusters naranjas distribuidos por Zaragoza, con la sidebar de fotos visible")
footer(s, 3)


# ========== SLIDE 4: ELEMENTO DIFERENCIADOR ==========
s = prs.slides.add_slide(BLANK)
set_bg(s, FONDO)
add_accent_bar(s, Inches(0.5), Inches(0.6))
add_text(s, Inches(0.75), Inches(0.55), Inches(12), Inches(0.7),
         "Elemento diferenciador", 32, color=AZUL, bold=True)
add_text(s, Inches(0.75), Inches(1.3), Inches(12), Inches(0.5),
         "Integración única de fuentes abiertas", 18, color=NARANJA, bold=True)

# Grid 2x2 con las capas diferenciales
cells = [
    ("🗺️", "Cartografía IGN", "Vuelos 1956, MTN 1915"),
    ("🏛️", "Catastro", "35.418 edificios con año"),
    ("📖", "Wikipedia", "Contexto enciclopédico"),
    ("🏛️", "Contexto histórico", "Alcaldes, eventos, noticias"),
]
for i, (icon, title, desc) in enumerate(cells):
    col, row = i % 2, i // 2
    x = Inches(1.2 + col * 5.8)
    y = Inches(2.3 + row * 2.2)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.3), Inches(1.9))
    box.fill.solid()
    box.fill.fore_color.rgb = BLANCO
    box.line.color.rgb = NARANJA
    box.line.width = Pt(1.5)
    add_text(s, x + Inches(0.3), y + Inches(0.4), Inches(0.8), Inches(1),
             icon, 32)
    add_text(s, x + Inches(1.3), y + Inches(0.35), Inches(3.8), Inches(0.5),
             title, 18, color=AZUL, bold=True)
    add_text(s, x + Inches(1.3), y + Inches(0.9), Inches(3.8), Inches(0.9),
             desc, 13, color=TEXTO_MUTE)
screenshot_hint(s, "Captura recomendada: 4 mini-capturas en mosaico -> capa PNOA 1956, capa Catastro coloreada por décadas, popup Wikipedia, panel Contexto histórico")
footer(s, 4)


# ========== SLIDE 5: ARQUITECTURA (CÓMO) ==========
s = prs.slides.add_slide(BLANK)
set_bg(s, FONDO)
add_accent_bar(s, Inches(0.5), Inches(0.6))
add_text(s, Inches(0.75), Inches(0.55), Inches(12), Inches(0.7),
         "Arquitectura", 32, color=AZUL, bold=True)
add_text(s, Inches(0.75), Inches(1.3), Inches(12), Inches(0.5),
         "Cliente  →  API  →  Base de datos espacial", 16, color=TEXTO_MUTE)

# Tres cajas horizontales conectadas con flechas
layers = [
    ("Frontend", "React + Vite", "Leaflet + TypeScript", AZUL),
    ("Backend", "FastAPI (Python)", "PostGIS + pyproj", NARANJA),
    ("Datos", "PostgreSQL + PostGIS", "Supabase + Storage", AZUL),
]
for i, (title, l1, l2, color) in enumerate(layers):
    x = Inches(0.8 + i * 4.3)
    y = Inches(2.8)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(3))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    add_text(s, x, y + Inches(0.5), Inches(3.7), Inches(0.6),
             title, 22, color=BLANCO, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.4), Inches(3.7), Inches(0.5),
             l1, 15, color=BLANCO, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(2.0), Inches(3.7), Inches(0.5),
             l2, 15, color=BLANCO, align=PP_ALIGN.CENTER)
    # Flecha entre cajas
    if i < 2:
        arrow_x = x + Inches(3.7) + Emu(20000)
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, Inches(4.0), Inches(0.5), Inches(0.6))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = TEXTO_MUTE
        arrow.line.fill.background()
screenshot_hint(s, "Captura recomendada: diagrama propio de arquitectura o pestaña Network de DevTools mostrando las llamadas /api/map, /api/photos, /api/buildings")
footer(s, 5)


# ========== SLIDE 6: DECISIONES TÉCNICAS (POR QUÉ) ==========
s = prs.slides.add_slide(BLANK)
set_bg(s, FONDO)
add_accent_bar(s, Inches(0.5), Inches(0.6))
add_text(s, Inches(0.75), Inches(0.55), Inches(12), Inches(0.7),
         "Decisiones técnicas", 32, color=AZUL, bold=True)
add_text(s, Inches(0.75), Inches(1.3), Inches(12), Inches(0.5),
         "Por qué cada elección", 18, color=NARANJA)

decisions = [
    ("WebP", "−65% tamaño", "Carga rápida móvil"),
    ("Clustering", "5.214 markers", "Rendimiento Leaflet"),
    ("PostGIS", "Índices espaciales", "Queries bbox <50ms"),
    ("Pooler IPv4", "Render free tier", "Sin IPv6 outbound"),
]
for i, (what, how, why) in enumerate(decisions):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(2.4 + row * 2.2)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.8), Inches(1.9))
    box.fill.solid()
    box.fill.fore_color.rgb = BLANCO
    box.line.color.rgb = AZUL_CLARO
    box.line.width = Pt(1)
    add_text(s, x + Inches(0.3), y + Inches(0.3), Inches(5), Inches(0.6),
             what, 22, color=AZUL, bold=True)
    add_text(s, x + Inches(0.3), y + Inches(0.95), Inches(5), Inches(0.5),
             how, 14, color=NARANJA, bold=True)
    add_text(s, x + Inches(0.3), y + Inches(1.35), Inches(5), Inches(0.5),
             why, 13, color=TEXTO_MUTE)
screenshot_hint(s, "Captura recomendada: zoom alto del mapa con clusters expandidos mostrando decenas de marcadores agrupados (demuestra el clustering)")
footer(s, 6)


# ========== SLIDE 7: DATOS ==========
s = prs.slides.add_slide(BLANK)
set_bg(s, FONDO)
add_accent_bar(s, Inches(0.5), Inches(0.6))
add_text(s, Inches(0.75), Inches(0.55), Inches(12), Inches(0.7),
         "Ingesta de datos (ETL)", 32, color=AZUL, bold=True)
add_text(s, Inches(0.75), Inches(1.3), Inches(12), Inches(0.5),
         "Pipeline de procesamiento", 18, color=TEXTO_MUTE)

# Flujo horizontal con 4 pasos
steps = [
    ("Origen", "GML + CSV"),
    ("Limpieza", "CP850 → UTF-8"),
    ("Transformar", "EPSG 25830 → 4326"),
    ("Cargar", "PostGIS + Storage"),
]
for i, (title, detail) in enumerate(steps):
    x = Inches(0.6 + i * 3.2)
    y = Inches(3.0)
    # Círculo con número
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.1), y, Inches(0.9), Inches(0.9))
    circle.fill.solid()
    circle.fill.fore_color.rgb = NARANJA
    circle.line.fill.background()
    add_text(s, x + Inches(1.1), y + Inches(0.15), Inches(0.9), Inches(0.6),
             str(i + 1), 28, color=BLANCO, bold=True, align=PP_ALIGN.CENTER)
    # Título y detalle
    add_text(s, x, y + Inches(1.2), Inches(3.1), Inches(0.5),
             title, 18, color=AZUL, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.75), Inches(3.1), Inches(0.5),
             detail, 13, color=TEXTO_MUTE, align=PP_ALIGN.CENTER)
    # Flecha
    if i < 3:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                    x + Inches(2.2), y + Inches(0.35),
                                    Inches(0.8), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = TEXTO_MUTE
        arrow.line.fill.background()

# Banner inferior con volumen
banner = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.1))
banner.fill.solid()
banner.fill.fore_color.rgb = AZUL
banner.line.fill.background()
add_text(s, Inches(0.8), Inches(5.75), Inches(11.7), Inches(0.6),
         "5.214 fotos   +   35.418 edificios   +   125 años de contexto histórico",
         18, color=BLANCO, bold=True, align=PP_ALIGN.CENTER)
screenshot_hint(s, "Captura recomendada: fragmento de código de fix_encoding_inplace.py o tabla SQL con las tildes correctas después del fix")
footer(s, 7)


# ========== SLIDE 8: DEMO / DESPLIEGUE ==========
s = prs.slides.add_slide(BLANK)
set_bg(s, FONDO)
add_accent_bar(s, Inches(0.5), Inches(0.6))
add_text(s, Inches(0.75), Inches(0.55), Inches(12), Inches(0.7),
         "Demo en producción", 32, color=AZUL, bold=True)
add_text(s, Inches(0.75), Inches(1.3), Inches(12), Inches(0.5),
         "Desplegado en Render + Supabase", 18, color=TEXTO_MUTE)

# Caja grande central con URL
url_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(1.5), Inches(2.5), Inches(10.3), Inches(3.5))
url_box.fill.solid()
url_box.fill.fore_color.rgb = BLANCO
url_box.line.color.rgb = AZUL
url_box.line.width = Pt(2)
add_text(s, Inches(1.5), Inches(3.0), Inches(10.3), Inches(0.6),
         "🌐", 40, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(3.9), Inches(10.3), Inches(0.6),
         "zaragoza-historica-frontend-static.onrender.com",
         20, color=AZUL, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(4.7), Inches(10.3), Inches(0.5),
         "Acceso público  ·  PWA  ·  Responsive móvil", 14, color=TEXTO_MUTE, align=PP_ALIGN.CENTER)
screenshot_hint(s, "Captura recomendada: montaje desktop + móvil de la web en producción, con una foto histórica abierta en fullscreen")
footer(s, 8)


# ========== SLIDE 9: REFLEXIÓN ==========
s = prs.slides.add_slide(BLANK)
set_bg(s, FONDO)
add_accent_bar(s, Inches(0.5), Inches(0.6))
add_text(s, Inches(0.75), Inches(0.55), Inches(12), Inches(0.7),
         "Reflexión", 32, color=AZUL, bold=True)
add_text(s, Inches(0.75), Inches(1.3), Inches(12), Inches(0.5),
         "Aprendizajes del proceso", 18, color=TEXTO_MUTE)

learnings = [
    ("✓", "Arquitectura full-stack real"),
    ("✓", "Datos espaciales (PostGIS)"),
    ("✓", "Despliegue cloud (Render + Supabase)"),
    ("✓", "Ingesta ETL de fuentes heterogéneas"),
    ("✓", "Optimización rendimiento web"),
]
for i, (check, text) in enumerate(learnings):
    y = Inches(2.4 + i * 0.75)
    # Check naranja
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), y, Inches(0.5), Inches(0.5))
    c.fill.solid()
    c.fill.fore_color.rgb = NARANJA
    c.line.fill.background()
    add_text(s, Inches(1.5), y + Inches(0.05), Inches(0.5), Inches(0.5),
             check, 18, color=BLANCO, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2.3), y + Inches(0.05), Inches(10), Inches(0.5),
             text, 18, color=TEXTO)
screenshot_hint(s, "Captura recomendada: dashboard de Render y Supabase mostrando los servicios desplegados y la BD con las tablas")
footer(s, 9)


# ========== SLIDE 10: CIERRE ==========
s = prs.slides.add_slide(BLANK)
set_bg(s, AZUL)
# Logo centrado
if LOGO_PATH.exists():
    s.shapes.add_picture(str(LOGO_PATH), Inches(5.67), Inches(1.5), width=Inches(2), height=Inches(2))
# Gracias
add_text(s, Inches(0.5), Inches(3.9), Inches(12.3), Inches(1.2),
         "Gracias", 60, color=BLANCO, bold=True, align=PP_ALIGN.CENTER)
# Separador naranja
sep = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.17), Inches(5.1), Inches(1), Emu(15000))
sep.fill.solid()
sep.fill.fore_color.rgb = NARANJA
sep.line.fill.background()
add_text(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.5),
         "¿Preguntas?", 22, color=BLANCO, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5),
         "zaragoza-historica-frontend-static.onrender.com", 13, color=NARANJA, align=PP_ALIGN.CENTER)


prs.save(str(OUT_PATH))
import sys
sys.stdout.reconfigure(encoding='utf-8')
print(f"OK -> {OUT_PATH}")
print(f"Slides: {len(prs.slides)}")
